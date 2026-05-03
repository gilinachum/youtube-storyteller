"""PowerPoint extraction tool for StoryTeller — supports local files and S3."""

import os
import tempfile
import boto3
from pptx import Presentation
from strands import tool

UPLOAD_BUCKET = os.environ.get("UPLOAD_BUCKET", "")
s3 = boto3.client("s3", region_name=os.environ.get("AWS_REGION", "us-west-2"))


def _resolve_file(file_path: str) -> str:
    """Download from S3 if path starts with s3:// or uploads/, otherwise return as-is."""
    s3_key = None

    if file_path.startswith("s3://"):
        path = file_path[5:]
        if path.startswith(UPLOAD_BUCKET + "/"):
            s3_key = path[len(UPLOAD_BUCKET) + 1:]
        elif path.startswith("uploads/"):
            s3_key = path
        else:
            parts = path.split("/", 1)
            s3_key = parts[1] if len(parts) > 1 else path

    elif file_path.startswith("uploads/"):
        s3_key = file_path

    if s3_key and UPLOAD_BUCKET:
        suffix = os.path.splitext(s3_key)[1] or ".pptx"
        tmp = tempfile.NamedTemporaryFile(suffix=suffix, delete=False)
        try:
            s3.download_file(UPLOAD_BUCKET, s3_key, tmp.name)
            return tmp.name
        except Exception as e:
            raise FileNotFoundError(f"Failed to download s3://{UPLOAD_BUCKET}/{s3_key}: {e}")

    return file_path


@tool
def pptx_extract(file_path: str) -> str:
    """Extract text and speaker notes from a PowerPoint file.

    Use this tool when the user provides a PPTX file or uploads one.
    Supports local paths and S3 references (s3://... or uploads/...).

    Args:
        file_path: Path to the PPTX file — local path or S3 key.

    Returns:
        Structured text with slide content and speaker notes per slide.
    """
    local_path = None
    try:
        local_path = _resolve_file(file_path)
        prs = Presentation(local_path)
        slides_text = []

        for i, slide in enumerate(prs.slides, 1):
            slide_content = []
            slide_content.append(f"=== Slide {i} ===")

            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)

            if texts:
                slide_content.append("Content:")
                for t in texts:
                    slide_content.append(f"  - {t}")

            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_content.append(f"Speaker Notes: {notes}")

            slides_text.append("\n".join(slide_content))

        if not slides_text:
            return f"[Warning] No content found in {file_path}"

        return f"PowerPoint content ({len(slides_text)} slides):\n\n" + "\n\n".join(slides_text)

    except FileNotFoundError as e:
        return f"[Error] PPTX file not found: {e}"
    except Exception as e:
        return f"[Error] Failed to extract PPTX: {str(e)}"
    finally:
        if local_path and local_path != file_path and os.path.exists(local_path):
            os.unlink(local_path)
